"""
PPTX Loader — python-pptx primary (richest slide structure), Docling fallback.

PPTX is a structured format, and python-pptx exposes the font-size / placeholder /
indent-level signals that recover what matters on a slide:
  * Slide TITLE  → heading      (title placeholder, else the largest-font text)
  * Bullets      → list_item    (paragraph indent level or a leading bullet char)
  * Body text    → paragraph
  * Tables       → table        (headers + rows)
  * Pictures     → image        (saved to disk, summarized later by the vision step)

Docling's PPTX backend discards those signals and flattens every slide into plain
paragraphs, so it is used only as a FALLBACK — if python-pptx is unavailable or
returns nothing usable, `_load_docling` runs the same shared
`docling_to_parsed_document()` pipeline as PDF/DOCX (one section per slide via
`split_on_page`). Either way the output is the same canonical element schema.
"""
import os
import re
import logging

logger = logging.getLogger(__name__)

_BULLET = re.compile(r"^\s*(?:[-—–•●▪◦*]|\d+\s*[.)])\s*\S")

_converter = None


def _get_converter():
    """Reuse one default Docling converter (same rationale as the DOCX loader:
    the default pipeline carries the option object convert() expects and its
    backend extracts embedded images — a bare custom PipelineOptions breaks it)."""
    global _converter
    if _converter is None:
        from docling.document_converter import DocumentConverter
        _converter = DocumentConverter()
        logger.info("[PPTX_LOADER] Docling converter created (default pipeline)")
    return _converter


def load(file_path: str) -> dict:
    """python-pptx first (richest slide structure); fall back to Docling."""
    try:
        result = _load_native(file_path)
        if result and result.get("raw_text", "").strip():
            return result
        logger.warning("[PPTX_LOADER] python-pptx produced no text; using Docling")
    except Exception as e:
        logger.warning(f"[PPTX_LOADER] python-pptx failed ({e}); using Docling")
    return _load_docling(file_path)


# ── Docling path (fallback — same shared pipeline as PDF/DOCX) ──

def _load_docling(file_path: str) -> dict:
    logger.info(f"[PPTX_LOADER] Loading with Docling: {os.path.basename(file_path)}")

    from app.services.providers.parsers._docling import docling_to_parsed_document
    from app.services.providers.loaders._utils import build_doc_metadata, clean_text
    from app.services.providers.parsers.hierarchy import (
        build_hierarchy, build_section_hierarchy,
    )

    converter = _get_converter()
    result = converter.convert(file_path)

    # Raw text from markdown export
    raw_text = clean_text(result.document.export_to_markdown())

    # ParsedDocument from structure — pass the real file_path so _docling saves
    # slide images under uploads/{space_id}/images (the folder the route serves).
    parsed_doc = docling_to_parsed_document(
        result=result,
        file_type="PPTX",
        category="document",
        metadata={"source": os.path.basename(file_path),
                  "file_path": os.path.abspath(file_path)},
    )

    num_pages = parsed_doc.num_pages or 1
    metadata = build_doc_metadata(file_path, num_pages, "pptx", parser_name="docling")

    build_hierarchy(parsed_doc.elements)
    build_section_hierarchy(parsed_doc.sections)
    parsed_doc.metadata = metadata

    logger.info(f"[PPTX_LOADER] Docling → {parsed_doc.total_sections} sections, "
                f"{parsed_doc.total_images} images, "
                f"{len(parsed_doc.elements)} elements")

    return {
        "raw_text": raw_text,
        "num_pages": num_pages,
        "file_type": "PPTX",
        "category": "document",
        "metadata": metadata,
        "total_chars": len(raw_text),
        "parsed_document": parsed_doc.to_dict(),
    }


# ── Native python-pptx path (primary) ──

def _load_native(file_path: str) -> dict:
    logger.info(f"[PPTX_LOADER] Parsing with python-pptx: {os.path.basename(file_path)}")

    from pptx import Presentation
    from app.services.providers.parsers.parsed_document import ParsedDocument, Section
    from app.services.providers.parsers._docling import _get_images_dir
    from app.services.providers.parsers.hierarchy import (
        build_hierarchy, build_section_hierarchy,
    )
    from app.services.providers.loaders._utils import build_doc_metadata, clean_text

    prs = Presentation(file_path)
    images_dir = _get_images_dir(os.path.abspath(file_path))

    sections = []
    elements = []
    all_images = []
    raw_parts = []
    ro = 0

    def _emit(el_type, content, page, level=None, list_level=None):
        nonlocal ro
        ro += 1
        el = {
            "id": f"e{ro}",
            "type": el_type,
            "content": content,
            "location": {"page": page, "bbox": []},
            "hierarchy": {"parent": None, "parent_section": None},
            "metadata": {"reading_order": ro},
        }
        if level is not None:
            el["level"] = level
        if list_level:
            el["metadata"]["list_level"] = list_level
        elements.append(el)

    for page, slide in enumerate(prs.slides, 1):
        texts, tables = [], []
        _walk(slide.shapes, texts, tables)

        title_shape = _pick_title(slide, texts)
        body = sorted([s for s in texts if s is not title_shape], key=_pos)

        slide_lines = []
        heading = ""

        if title_shape is not None:
            heading = title_shape.text_frame.text.strip()
            if heading:
                _emit("heading", {"text": heading}, page, level=1)
                slide_lines.append(heading)

        for shape in body:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if not t:
                    continue
                slide_lines.append(t)
                lvl = getattr(para, "level", 0) or 0
                if lvl > 0 or _BULLET.match(t):
                    _emit("list_item", {"text": t}, page, list_level=lvl + 1)
                else:
                    _emit("paragraph", {"text": t}, page)

        # Tables (python-pptx exposes them directly)
        for tshape in tables:
            headers, rows, md = _table_data(tshape)
            if headers or rows:
                _emit("table", {"caption": "", "headers": headers,
                                "rows": rows, "markdown": md}, page)
                slide_lines.append(md)

        # Pictures → save + emit image elements. Pulled from the slide's
        # relationships (not the shape tree) so it catches picture shapes AND
        # picture-fills / freeform-fills / backgrounds — which python-pptx does
        # not expose as PICTURE shapes and _save_pic used to miss entirely.
        for idx, (blob, ext) in enumerate(_slide_images(slide), 1):
            saved = _save_blob(blob, ext, images_dir, page, idx)
            if saved:
                all_images.append(saved)
                _emit("image", {
                    "image_path": saved.image_path, "caption": "Image",
                    "text_for_embedding": "", "ocr_text": "",
                }, page)

        body_content = "\n".join(l for l in slide_lines if l != heading)
        sections.append(Section(heading=heading, content=body_content, level=1, page=page))
        if slide_lines:
            raw_parts.append("\n".join(slide_lines))

    raw_text = clean_text("\n\n".join(raw_parts))
    if not raw_text.strip() and not sections:
        raise ValueError("PPTX contains no readable text")

    num_slides = len(prs.slides)
    metadata = build_doc_metadata(file_path, num_slides, "pptx", parser_name="python-pptx")

    parsed_doc = ParsedDocument(
        title=sections[0].heading if sections and sections[0].heading else "",
        sections=sections,
        images=all_images,
        elements=elements,
        metadata=metadata,
        num_pages=num_slides,
        file_type="PPTX",
        category="document",
    )

    build_hierarchy(parsed_doc.elements)
    build_section_hierarchy(parsed_doc.sections)

    logger.info(f"[PPTX_LOADER] {num_slides} slide(s) → "
                f"{parsed_doc.total_sections} sections, "
                f"{parsed_doc.total_images} images, "
                f"{len(parsed_doc.elements)} elements")

    return {
        "raw_text": raw_text,
        "num_pages": num_slides,
        "file_type": "PPTX",
        "category": "document",
        "metadata": metadata,
        "total_chars": len(raw_text),
        "parsed_document": parsed_doc.to_dict(),
    }


# ── shape helpers (native path) ──

def _walk(shapes, texts, tables):
    """Recurse into groups, bucketing shapes into text / table. Pictures are
    handled separately from the slide's relationships (see _slide_images), which
    catches images the shape tree hides (fills, freeforms, backgrounds)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for s in shapes:
        try:
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                _walk(s.shapes, texts, tables)
                continue
        except Exception:
            pass
        try:
            if s.has_table:
                tables.append(s)
                continue
        except Exception:
            pass
        try:
            if s.has_text_frame and s.text_frame.text.strip():
                texts.append(s)
        except Exception:
            pass


def _pos(shape):
    try:
        return (shape.top or 0, shape.left or 0)
    except Exception:
        return (0, 0)


def _max_font(shape) -> float:
    """Largest run font size (pt) in a shape, 0 if unknown."""
    best = 0.0
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                sz = getattr(run.font, "size", None)
                if sz is not None:
                    best = max(best, float(sz.pt))
    except Exception:
        pass
    return best


def _pick_title(slide, texts):
    """Slide title: a real TITLE placeholder, else the largest-font short text."""
    for s in texts:
        try:
            if s.is_placeholder and "TITLE" in str(s.placeholder_format.type):
                return s
        except Exception:
            pass
    best, best_sz = None, 0.0
    for s in texts:
        txt = s.text_frame.text.strip()
        if not txt or len(txt) > 120:
            continue
        sz = _max_font(s)
        if sz > best_sz:
            best_sz, best = sz, s
    return best


def _table_data(shape):
    """Return (headers, rows, markdown) for a python-pptx table shape."""
    try:
        tbl = shape.table
    except Exception:
        return [], [], ""
    grid = [[c.text.strip() for c in row.cells] for row in tbl.rows]
    if not grid:
        return [], [], ""
    headers = grid[0]
    rows = [
        {headers[i] if i < len(headers) else f"col{i}": v for i, v in enumerate(r)}
        for r in grid[1:]
    ]
    md_lines = [" | ".join(headers), " | ".join("---" for _ in headers)]
    md_lines += [" | ".join(r) for r in grid[1:]]
    return headers, rows, "\n".join(md_lines)


def _slide_images(slide):
    """Every embedded image a slide references, via its part relationships.

    A .pptx often places images as shape/slide FILLS or freeform fills rather
    than as PICTURE shapes, so walking the shape tree misses them (python-pptx
    reports 0 pictures even when ppt/media/ holds several PNGs). The slide's
    relationships list every image it embeds, regardless of how it's used.

    Returns a list of (blob, ext), de-duplicated per slide. External (linked)
    images and non-image parts are skipped.
    """
    out, seen = [], set()
    try:
        rels = slide.part.rels
    except Exception:
        return out
    for rel in rels.values():
        try:
            if getattr(rel, "is_external", False):
                continue
            if "image" not in (rel.reltype or "").lower():
                continue
            part = rel.target_part
            name = str(getattr(part, "partname", ""))
            if name in seen:
                continue
            seen.add(name)
            blob = part.blob
            if not blob:
                continue
            ext = name.rsplit(".", 1)[-1].lower() if "." in name else "png"
            if ext == "jpg":
                ext = "jpeg"
            out.append((blob, ext))
        except Exception:
            continue
    return out


def _save_blob(blob, ext, images_dir, page, idx):
    """Write an image blob to disk; return a ParsedDocument Image or None."""
    from app.services.providers.parsers.parsed_document import Image
    if not images_dir or not blob:
        return None
    try:
        os.makedirs(images_dir, exist_ok=True)
        path = os.path.join(images_dir, f"page{page}_img{idx}.{ext}")
        with open(path, "wb") as f:
            f.write(blob)
        return Image(caption="Image", ocr_text="", image_path=path, page=page, bbox=[])
    except Exception as e:
        logger.warning(f"[PPTX_LOADER] could not save one image: {e}")
        return None
